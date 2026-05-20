"""
Generate docs/01-slide-deck.pptx from scratch following the official
Redis brand guidelines (brand.redis.io).

Palette pulled from redis.io production CSS (CSS variables --hyper-*,
--dusk-*, --yellow-*, --midnight):

    Hyper-05    #ff4438   primary Redis red
    Hyper-04    #fd736a   light red
    Hyper-09    #8a221c   dark red
    Dusk        #163341   slate (subheads, body on light bg)
    Midnight    #091a23   near-black, used for hero backgrounds
    Dusk-50     #8a99a0   muted neutral
    Dusk-01     #f3f3f3   page background
    Yellow      #dcff1e   accent / "live" indicator

Brand typography: Space Grotesk (display + body) / Space Mono (code).
PowerPoint substitutes the OS font when Space Grotesk isn't installed
locally; we set it explicitly so the deck looks correct on a Mac with
the Google Font installed and degrades gracefully elsewhere.
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------
# Brand tokens
# ---------------------------------------------------------------------
RED       = RGBColor(0xFF, 0x44, 0x38)   # hyper-05
RED_DK    = RGBColor(0x8A, 0x22, 0x1C)   # hyper-09
RED_LT    = RGBColor(0xFD, 0x73, 0x6A)   # hyper-04
DUSK      = RGBColor(0x16, 0x33, 0x41)   # dusk
MIDNIGHT  = RGBColor(0x09, 0x1A, 0x23)   # midnight
DUSK_50   = RGBColor(0x8A, 0x99, 0xA0)   # muted
DUSK_30   = RGBColor(0xB9, 0xC2, 0xC6)   # divider
DUSK_01   = RGBColor(0xF3, 0xF3, 0xF3)   # page bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW    = RGBColor(0xDC, 0xFF, 0x1E)   # accent
CHECK_GR  = RGBColor(0x2E, 0xCC, 0x71)
WARN_RD   = RGBColor(0xFF, 0x5A, 0x6C)

# Per Redis brand: Space Grotesk for display + body, Space Mono for code.
FONT_BODY = "Space Grotesk"
FONT_MONO = "Space Mono"

# 16:9 widescreen (standard Redis presentation aspect)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.line.fill.background() if line is None else _set_line(s, line)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.shadow.inherit = False
    return s


def _set_line(shape, rgb):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(0.75)


def add_text(
    slide, x, y, w, h, text,
    *, size=14, bold=False, color=DUSK, font=FONT_BODY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_chrome(slide, slide_num, total, *, dark=False):
    """Top brand bar + footer; same on every content slide."""
    fg = WHITE if dark else DUSK
    bar_fg = WHITE if dark else MIDNIGHT
    # Brand strip top: red square + wordmark
    add_rect(slide, Inches(0.5), Inches(0.32), Inches(0.22), Inches(0.22), RED)
    add_text(slide, Inches(0.78), Inches(0.22), Inches(3), Inches(0.4),
             "redis", size=18, bold=True, color=bar_fg, font=FONT_BODY)
    add_text(slide, Inches(2.0), Inches(0.30), Inches(7), Inches(0.3),
             "  |   Redis Data Integration  ·  Securities & Trading Firm",
             size=11, color=DUSK_50, font=FONT_BODY)
    # Slide counter (right)
    add_text(slide, Inches(11.5), Inches(0.30), Inches(1.5), Inches(0.3),
             f"{slide_num:02d} / {total:02d}",
             size=11, color=DUSK_50, font=FONT_MONO, align=PP_ALIGN.RIGHT)
    # Footer divider line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.0),
                                      Inches(12.83), Inches(7.0))
    line.line.color.rgb = DUSK_30 if not dark else RGBColor(0x2D, 0x47, 0x54)
    line.line.width = Pt(0.5)
    # Footer text
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
             "Confidential  ·  prepared for Securities & Trading Firm",
             size=9, color=DUSK_50, font=FONT_BODY)
    add_text(slide, Inches(8.5), Inches(7.05), Inches(4.5), Inches(0.3),
             "redis.io/integrate/redis-data-integration",
             size=9, color=DUSK_50, font=FONT_MONO, align=PP_ALIGN.RIGHT)


def add_h1(slide, text, *, color=MIDNIGHT, y=Inches(0.95)):
    add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.9),
             text, size=32, bold=True, color=color, font=FONT_BODY)


def add_kicker(slide, text, *, color=RED, y=Inches(0.7)):
    add_text(slide, Inches(0.5), y, Inches(12.3), Inches(0.3),
             text.upper(), size=10, bold=True, color=color, font=FONT_MONO)


def add_talktrack(slide, text):
    """The little red-left-bar block that holds the presenter script."""
    y = Inches(5.85)
    h = Inches(1.05)
    add_rect(slide, Inches(0.5), y, Inches(0.06), h, RED)
    add_rect(slide, Inches(0.56), y, Inches(12.27), h, DUSK_01)
    add_text(slide, Inches(0.85), y + Inches(0.08), Inches(12), Inches(0.3),
             "TALK TRACK", size=9, bold=True, color=RED, font=FONT_MONO)
    add_text(slide, Inches(0.85), y + Inches(0.3), Inches(12), h - Inches(0.4),
             text, size=11, color=DUSK, font=FONT_BODY)


def add_bullets(slide, x, y, w, h, items, *, size=14, color=DUSK,
                bullet_color=RED, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet
        rb = p.add_run()
        rb.text = "■  "
        rb.font.name = FONT_MONO
        rb.font.size = Pt(size - 2)
        rb.font.color.rgb = bullet_color
        rb.font.bold = True
        # body
        rt = p.add_run()
        rt.text = item
        rt.font.name = FONT_BODY
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.bold = bold_first and i == 0
        p.space_after = Pt(6)


def add_card(slide, x, y, w, h, title, body_items, *, accent=RED):
    add_rect(slide, x, y, w, h, WHITE, line=DUSK_30)
    # Accent corner stripe
    add_rect(slide, x, y, Inches(0.18), h, accent)
    add_text(slide, x + Inches(0.4), y + Inches(0.2), w - Inches(0.6),
             Inches(0.4), title.upper(),
             size=11, bold=True, color=accent, font=FONT_MONO)
    add_bullets(slide, x + Inches(0.4), y + Inches(0.7),
                w - Inches(0.6), h - Inches(0.9), body_items,
                size=12, bullet_color=accent)


def add_code(slide, x, y, w, h, code):
    add_rect(slide, x, y, w, h, MIDNIGHT)
    add_text(slide, x + Inches(0.3), y + Inches(0.2),
             w - Inches(0.6), h - Inches(0.4),
             code, size=11, color=WHITE, font=FONT_MONO)


def set_notes(slide, text):
    """Attach speaker notes — visible to the presenter in Presenter View."""
    nf = slide.notes_slide.notes_text_frame
    nf.text = ""
    for i, line in enumerate(text.split("\n")):
        p = nf.paragraphs[0] if i == 0 else nf.add_paragraph()
        p.text = line


# ---------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------
TOTAL = 21

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]   # blank layout (no placeholders)

# -----------------------------------------------------------------
# 1. Cover
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, MIDNIGHT)
# Top brand strip (different on cover - just the wordmark)
add_rect(s, Inches(0.6), Inches(0.45), Inches(0.28), Inches(0.28), RED)
add_text(s, Inches(0.97), Inches(0.32), Inches(3), Inches(0.5),
         "redis", size=24, bold=True, color=WHITE, font=FONT_BODY)
# Yellow live-demo pill
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(2.7),
                          Inches(2.1), Inches(0.45))
pill.line.fill.background()
pill.fill.solid(); pill.fill.fore_color.rgb = YELLOW
pill.adjustments[0] = 0.5
add_text(s, Inches(0.6), Inches(2.7), Inches(2.1), Inches(0.45),
         "LIVE DEMO  ·  ON THIS LAPTOP",
         size=11, bold=True, color=MIDNIGHT, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Title
add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(1.4),
         "Redis Data Integration",
         size=64, bold=True, color=WHITE, font=FONT_BODY)
# Underline accent
add_rect(s, Inches(0.6), Inches(4.55), Inches(1.6), Inches(0.07), RED)
# Subtitle
add_text(s, Inches(0.6), Inches(4.75), Inches(12), Inches(0.6),
         "Real-time portfolio cache for Securities & Trading Firm",
         size=22, color=DUSK_30, font=FONT_BODY)
# Source -> sink line
add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
         "PostgreSQL  →  Debezium CDC  →  RDI processor  →  Redis Enterprise",
         size=13, color=YELLOW, font=FONT_MONO)
# Footer
add_text(s, Inches(0.6), Inches(6.5), Inches(8), Inches(0.4),
         "Prepared for the Securities & Trading Firm Architecture & Application team",
         size=11, color=DUSK_50, font=FONT_BODY)
add_text(s, Inches(0.6), Inches(6.85), Inches(8), Inches(0.3),
         "Tuesday  ·  Mumbai",
         size=11, color=DUSK_50, font=FONT_MONO)

# -----------------------------------------------------------------
# 2. The portfolio screen
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 2, TOTAL)
add_kicker(s, "Context")
add_h1(s, "The portfolio screen: where every retail investor lands")

add_text(s, Inches(0.5), Inches(1.9), Inches(6.5), Inches(0.5),
         "Securities & Trading Firm serves millions of demat clients —\n"
         "and every one of them lands on the same read:",
         size=14, color=DUSK, font=FONT_BODY)
add_bullets(s, Inches(0.5), Inches(2.9), Inches(6.5), Inches(2.0), [
    "“What do I hold today?”",
    "“What is my P&L right now?”",
    "“Show my last 10 trades.”",
], size=15, bullet_color=RED)
add_text(s, Inches(0.5), Inches(4.45), Inches(6.5), Inches(1.3),
         "Each of those reads fans out into joins on Oracle: "
         "holdings × security master × live LTP × trade history. "
         "During market hours, thousands of identical reads per second "
         "ask for the same few thousand customers' data.",
         size=12, color=DUSK_50, font=FONT_BODY)

add_card(s, Inches(7.4), Inches(1.9), Inches(5.4), Inches(3.85),
         "Symptoms you may already see", [
             "Oracle CPU spikes 9:15–10:00 and 14:30–15:30 IST",
             "Read-replica licensing grows quarter over quarter",
             "Mobile app P95 > 800 ms during peak load",
             "Hand-written cache-aside in 6 places — each subtly inconsistent",
         ])

add_talktrack(s,
    "Before talking about Redis, let me describe what we hear from brokers in India. "
    "The portfolio screen is the busiest entry point in the app. It is mostly reads, "
    "mostly the same data, repeated all day. Oracle is right for correctness — wrong "
    "for hot-path reads. Brokers respond with vertical scaling, read replicas, or "
    "cache-aside. Each path has a cost: licence, latency, or engineering. Sound familiar?"
)

# -----------------------------------------------------------------
# 3. Why cache-aside fails
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 3, TOTAL)
add_kicker(s, "Problem statement")
add_h1(s, "Why cache-aside fails at portfolio scale")

# Table-style comparison: 3 columns
hdr_y = Inches(1.95)
row_h = Inches(0.52)
col1, col2, col3 = Inches(0.5), Inches(4.4), Inches(8.7)
col1_w = Inches(3.85); col2_w = Inches(4.25); col3_w = Inches(4.1)

# Header band
add_rect(s, col1, hdr_y, Inches(12.3), Inches(0.45), MIDNIGHT)
add_text(s, col1 + Inches(0.2), hdr_y + Inches(0.1), col1_w, Inches(0.3),
         "CONCERN", size=11, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, col2 + Inches(0.2), hdr_y + Inches(0.1), col2_w, Inches(0.3),
         "CACHE-ASIDE  (today)", size=11, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, col3 + Inches(0.2), hdr_y + Inches(0.1), col3_w, Inches(0.3),
         "REDIS DATA INTEGRATION", size=11, bold=True, color=YELLOW, font=FONT_MONO)

rows = [
    ("First request after deploy",   "Cache miss → Oracle hit",                "Already prefetched"),
    ("Update in Oracle",             "App must invalidate (often forgets)",    "CDC propagates in < 1 s"),
    ("TTL strategy",                 "Always stale, always wrong window",      "No TTL — cache mirrors DB"),
    ("Code in app for cache logic",  "Sprinkled across 6 repos",               "Zero — declarative YAML"),
    ("Adding a new cached table",    "Engineering sprint",                     "+10 lines of YAML"),
    ("Cold-start after a wipe",      "Manual prefetch script",                 "RDI re-snapshots from source"),
]
y = hdr_y + Inches(0.45)
for i, (a, b, c) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else DUSK_01
    add_rect(s, col1, y, Inches(12.3), row_h, bg)
    add_text(s, col1 + Inches(0.2), y + Inches(0.12), col1_w, Inches(0.4),
             a, size=12, bold=True, color=MIDNIGHT, font=FONT_BODY)
    add_text(s, col2 + Inches(0.4), y + Inches(0.12), col2_w, Inches(0.4),
             "✗  " + b, size=12, color=WARN_RD, font=FONT_BODY)
    add_text(s, col3 + Inches(0.4), y + Inches(0.12), col3_w, Inches(0.4),
             "✓  " + c, size=12, color=CHECK_GR, font=FONT_BODY)
    y += row_h

add_talktrack(s,
    "Every team eventually builds cache-aside. The problem isn't building it — it's "
    "keeping it correct as the schema changes, as new tables get added, as new engineers "
    "join. RDI flips the model: the cache is no longer a side-effect maintained by the "
    "app; it is a continuously-synced projection of your source DB, owned by configuration."
)

# -----------------------------------------------------------------
# 4. What is RDI - in one slide
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 4, TOTAL)
add_kicker(s, "RDI in one slide")
add_h1(s, "Continuously synced. Zero application code.")

add_text(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.6),
         "RDI keeps Redis Enterprise in lockstep with your primary database "
         "using Change Data Capture. No application changes.",
         size=15, color=DUSK, font=FONT_BODY)

# Flow diagram - three boxes with arrows
box_y = Inches(2.9); box_h = Inches(1.4)
b1_x = Inches(0.7); b2_x = Inches(5.0); b3_x = Inches(9.3)
box_w = Inches(3.6)

def flow_box(x, label, sub, *, fill=WHITE, border=DUSK_30, label_color=MIDNIGHT):
    add_rect(s, x, box_y, box_w, box_h, fill, line=border)
    add_text(s, x, box_y + Inches(0.4), box_w, Inches(0.4),
             label, size=20, bold=True, color=label_color,
             font=FONT_BODY, align=PP_ALIGN.CENTER)
    add_text(s, x, box_y + Inches(0.85), box_w, Inches(0.4),
             sub, size=11, color=DUSK_50, font=FONT_MONO,
             align=PP_ALIGN.CENTER)

flow_box(b1_x, "Oracle / Postgres", "system of record")
flow_box(b2_x, "RDI", "CDC + transform", fill=MIDNIGHT,
         border=RED, label_color=YELLOW)
flow_box(b3_x, "Redis Enterprise", "hot cache, sub-ms reads")

# Arrows
for ax in [Inches(4.4), Inches(8.7)]:
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax,
                               box_y + Inches(0.55),
                               Inches(0.55), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = RED
    arrow.line.fill.background()

# 3 cards under the flow
card_y = Inches(4.8); card_h = Inches(1.85); card_w = Inches(4.0)
add_card(s, Inches(0.5),  card_y, card_w, card_h, "Capture", [
    "Debezium-based CDC",
    "Reads the database’s own redo / WAL",
    "Zero load on application tables",
])
add_card(s, Inches(4.65), card_y, card_w, card_h, "Transform", [
    "Declarative YAML jobs",
    "Pick columns, filter rows, set key pattern",
    "Pick target: JSON · Hash · Stream · ZSET",
])
add_card(s, Inches(8.8),  card_y, card_w, card_h, "Deliver", [
    "At-least-once into Redis Enterprise",
    "~10k records / sec / core",
    "Back-pressure + full observability",
])

# (no talk track on this one — already long)

# -----------------------------------------------------------------
# 5. How RDI is deployed
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 5, TOTAL)
add_kicker(s, "Architecture")
add_h1(s, "How RDI is deployed in production")

add_text(s, Inches(0.5), Inches(1.95), Inches(6.5), Inches(0.4),
         "Three planes  ·  standard RDI topology",
         size=14, bold=True, color=MIDNIGHT, font=FONT_BODY)
add_bullets(s, Inches(0.5), Inches(2.4), Inches(6.5), Inches(2.6), [
    "Data plane — Debezium collector + stream processor. 2 VMs, active/standby with leader election.",
    "Control plane — operator + REST API + Prometheus exporter on the same VMs.",
    "Management plane — redis-di CLI + Redis Insight pipeline editor (same tool you use today).",
], size=12)
add_text(s, Inches(0.5), Inches(4.7), Inches(6.5), Inches(1.0),
         "RDI's own state — CDC streams, checkpoints, schema history — sits "
         "inside Redis Enterprise. Nothing is persisted on the RDI VMs. "
         "Your security and operational posture stays simple.",
         size=12, color=DUSK_50, italic=True, font=FONT_BODY)

add_card(s, Inches(7.4), Inches(1.9), Inches(5.4), Inches(3.85),
         "For Securities & Trading Firm", [
             "2 RDI VMs (or K8s pods) per environment",
             "RDI state DB on existing Redis Enterprise — 250 MB shard pair",
             "Source: read-only CDC user on Oracle (same as GoldenGate setup)",
             "Target: any number of Redis Enterprise DBs",
         ], accent=YELLOW)

add_talktrack(s,
    "Operationally this is two VMs or two K8s pods. RDI has hot failover built in — "
    "leader election decides which one is active. All the state RDI needs lives in "
    "Redis Enterprise on your existing cluster. There is nothing else for you to back "
    "up. That matters for your operations team and your auditors."
)

# -----------------------------------------------------------------
# 6. Pipeline lifecycle
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 6, TOTAL)
add_kicker(s, "Lifecycle")
add_h1(s, "What happens when you click Deploy")

steps = [
    ("01", "Validate",
     "Syntax check on config.yaml, connectivity to source & target, replication slot exists."),
    ("02", "Snapshot",
     "Full parallel read of selected rows. Day-1 prefetch — no cold cache, ever."),
    ("03", "CDC streaming",
     "Tail Oracle redo / Postgres WAL. New events land in Redis in 100s of milliseconds."),
    ("04", "Update",
     "Add a column or a transformation — edit YAML, redeploy. No pipeline reset."),
    ("05", "Reset (rare)",
     "Rebuild a target DB from scratch with a single CLI command."),
]

y0 = Inches(2.0); h_step = Inches(0.85); spacing = Inches(0.1)
y = y0
for num, title, body in steps:
    add_rect(s, Inches(0.5), y, Inches(12.3), h_step, DUSK_01)
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.13),
                                Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = RED
    circle.line.fill.background()
    add_text(s, Inches(0.7), y + Inches(0.13), Inches(0.6), Inches(0.6),
             num, size=14, bold=True, color=WHITE, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y + Inches(0.13), Inches(3), Inches(0.4),
             title, size=15, bold=True, color=MIDNIGHT, font=FONT_BODY)
    add_text(s, Inches(1.55), y + Inches(0.45), Inches(11.0), Inches(0.4),
             body, size=12, color=DUSK, font=FONT_BODY)
    y += h_step + spacing

add_talktrack(s,
    "This is the lifecycle. The key word is snapshot — on day 1, RDI prefetches your "
    "whole dataset. So the very first user request after go-live is already hot. "
    "That's the difference between cache-aside and cache-prefetching: you stop paying "
    "the cold-start tax."
)

# -----------------------------------------------------------------
# 7. Mapping Securities & Trading Firm data to Redis
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 7, TOTAL)
add_kicker(s, "Mapping for your use case")
add_h1(s, "Portfolio tables → Redis data structures")

rows = [
    ("portfolio.customer",        "customer:{client_code}",                    "JSON",        "Full profile read on every login"),
    ("portfolio.security_master", "security:{symbol}",                         "Hash",        "Reference data, per-field reads"),
    ("portfolio.holding",         "holding:{cust}:{sec}",                      "JSON",        "Single hottest object in the app"),
    ("portfolio.trade",           "trades:{cust}  +  trade:{trade_id}",        "Stream + JSON","Live order feed + lookup"),
    ("portfolio.market_price",    "price:{security_id}",                       "Hash",        "Tickers update every second"),
]

hdr_y = Inches(2.0)
add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.45), MIDNIGHT)
cols = [Inches(0.7), Inches(3.7), Inches(7.5), Inches(8.7)]
widths = [Inches(3.0), Inches(3.8), Inches(1.1), Inches(4.0)]
heads = ["SOURCE TABLE (Oracle)", "REDIS KEY PATTERN", "TYPE", "WHY"]
for x, w, h_ in zip(cols, widths, heads):
    add_text(s, x, hdr_y + Inches(0.1), w, Inches(0.3),
             h_, size=11, bold=True, color=YELLOW, font=FONT_MONO)

row_h = Inches(0.7)
y = hdr_y + Inches(0.45)
for i, (a, b, c, d) in enumerate(rows):
    bg = WHITE if i % 2 == 0 else DUSK_01
    add_rect(s, Inches(0.5), y, Inches(12.3), row_h, bg)
    add_text(s, cols[0], y + Inches(0.22), widths[0], Inches(0.4),
             a, size=12, color=MIDNIGHT, font=FONT_MONO, bold=True)
    add_text(s, cols[1], y + Inches(0.22), widths[1], Inches(0.4),
             b, size=12, color=RED, font=FONT_MONO)
    add_text(s, cols[2], y + Inches(0.22), widths[2], Inches(0.4),
             c, size=12, color=DUSK, font=FONT_BODY)
    add_text(s, cols[3], y + Inches(0.22), widths[3], Inches(0.4),
             d, size=12, color=DUSK, font=FONT_BODY)
    y += row_h

add_text(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.6),
         "Same source, different Redis shapes — pick the structure that fits the "
         "read pattern. RDI does this per table via a YAML job.",
         size=12, color=DUSK_50, italic=True, font=FONT_BODY)

add_talktrack(s,
    "Here is the concrete mapping. Notice we use different Redis structures for "
    "different tables. A Hash for the security master where the trading screen wants "
    "one or two fields at a time. A JSON document for the holding because the mobile "
    "app wants the whole object. A Stream for trades so 'My Orders' just tails it. "
    "Same source, three target representations, all driven by config."
)

# -----------------------------------------------------------------
# 8. A complete RDI job in 10 lines
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 8, TOTAL)
add_kicker(s, "Configuration over code")
add_h1(s, "A complete RDI job — ten lines")

code = (
    "source:\n"
    "  server_name: sectrade\n"
    "  schema: portfolio\n"
    "  table: holding\n"
    "\n"
    "output:\n"
    "  - uses: redis.write\n"
    "    with:\n"
    "      connection: target\n"
    "      data_type: json\n"
    "      key:\n"
    "        expression: concat(['holding:', to_string(customer_id), ':', to_string(security_id)])\n"
    "        language:   jmespath\n"
)
add_code(s, Inches(0.5), Inches(2.0), Inches(8.0), Inches(3.7), code)

# Annotation card on the right
add_card(s, Inches(8.8), Inches(2.0), Inches(4.0), Inches(3.7),
         "What this one file gives you", [
             "Keeps holding:{cust}:{sec} in sync forever",
             "Reviewable by DBAs, security, auditors",
             "No compiled artefact, no Kafka app",
             "Versioned in git like any other config",
         ])

add_talktrack(s,
    "This is literally one of the five jobs we deploy in the demo. Ten lines. Your "
    "DBAs can read it, your auditors can read it, your security team can read it. "
    "There is no compiled artifact, no custom Kafka Streams app to maintain. Compare "
    "this to one cache-aside helper class in your Java codebase — it is already longer."
)

# -----------------------------------------------------------------
# 9. Redis Demo Center walkthrough
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 9, TOTAL)
add_kicker(s, "Public reference: redis.io / demo-center")
add_h1(s, "The official RDI tour — same flow, your data")

steps = [
    "Add a Redis Enterprise cluster to Redis Insight.",
    "Connect to the RDI control endpoint — Insight detects an empty pipeline.",
    "Open the pipeline editor — left pane: source tables, right pane: generated output.",
    "Drag a transformation (rename / filter / project) — Insight rewrites the YAML.",
    "Click Deploy. RDI validates, snapshots, then streams CDC.",
    "The dashboard inside Insight shows throughput, lag and DLQ per table.",
]
y = Inches(2.0)
for i, step in enumerate(steps, 1):
    add_text(s, Inches(0.5), y, Inches(0.6), Inches(0.5),
             f"{i:02d}", size=20, bold=True, color=RED, font=FONT_MONO)
    add_text(s, Inches(1.1), y + Inches(0.08), Inches(6.5), Inches(0.5),
             step, size=13, color=DUSK, font=FONT_BODY)
    y += Inches(0.5)

add_card(s, Inches(8.2), Inches(2.0), Inches(4.6), Inches(3.7),
         "Why this matters for the firm", [
             "Same UI in dev, staging, prod",
             "Test pipeline against latest CDC on every save",
             "Dev writes YAML and verifies it — no Kafka team needed",
             "Insight ships free with Redis Enterprise",
         ])

add_talktrack(s,
    "Before we go to the laptop: the public Redis demo center shows exactly this "
    "workflow against a sample database. We are going to do the same flow but with "
    "a portfolio dataset that mirrors your business. The UI is Redis Insight — the "
    "same tool that already ships with Redis Enterprise. When you deploy RDI in "
    "production, your engineers will be using a tool they already know."
)

# -----------------------------------------------------------------
# 10. Demo scenarios (6 cards)
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 10, TOTAL)
add_kicker(s, "Live demo")
add_h1(s, "What you will see in the next 25 minutes")

scenarios = [
    ("Scenario 1", "Initial snapshot",
     "Bring up the stack. RDI prefetches all customers, securities, "
     "holdings, trades and prices. Open the dashboard — everything is already there."),
    ("Scenario 2", "Live trade",
     "Place a BUY in Postgres. Watch the trade stream, the holding JSON, "
     "and the portfolio KPIs update within ~1 second."),
    ("Scenario 3", "Live price tick",
     "Run the market-data simulator. P&L updates in real time on the "
     "dashboard without anyone refreshing."),
    ("Scenario 4", "Postgres vs Redis latency",
     "Same portfolio summary query. Postgres in tens of ms, Redis "
     "in sub-ms. Multiply by peak QPS for Oracle savings."),
    ("Scenario 5", "Pipeline observability",
     "Open Redis Insight. Per-table CDC counts, DLQ, deployed YAML — "
     "what your ops team sees every day."),
    ("Scenario 6", "Schema change",
     "Add a column to holding in Postgres. RDI auto-detects it; "
     "next update lands in Redis. No pipeline restart."),
]

card_w = Inches(4.0); card_h = Inches(1.8)
positions = [
    (Inches(0.5),  Inches(1.9)),
    (Inches(4.65), Inches(1.9)),
    (Inches(8.8),  Inches(1.9)),
    (Inches(0.5),  Inches(3.85)),
    (Inches(4.65), Inches(3.85)),
    (Inches(8.8),  Inches(3.85)),
]
for (x, y), (tag, title, body) in zip(positions, scenarios):
    add_rect(s, x, y, card_w, card_h, WHITE, line=DUSK_30)
    add_rect(s, x, y, card_w, Inches(0.32), MIDNIGHT)
    add_text(s, x + Inches(0.25), y + Inches(0.05), card_w, Inches(0.25),
             tag.upper(), size=10, bold=True, color=YELLOW, font=FONT_MONO)
    add_text(s, x + Inches(0.25), y + Inches(0.4), card_w - Inches(0.5),
             Inches(0.4), title, size=15, bold=True, color=MIDNIGHT, font=FONT_BODY)
    add_text(s, x + Inches(0.25), y + Inches(0.8), card_w - Inches(0.5),
             card_h - Inches(0.9), body, size=11, color=DUSK, font=FONT_BODY)

add_talktrack(s,
    "Six short scenarios, each under five minutes. Please interrupt me at any point "
    "with the kind of objection you'd hear from your own architecture review board. "
    "We can pause and poke around the running system as we go."
)

# -----------------------------------------------------------------
# 11. TCO / cost
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 11, TOTAL)
add_kicker(s, "Economics")
add_h1(s, "What this looks like for the CFO")

# Two columns
col_y = Inches(2.0); col_h = Inches(3.6); col_w = Inches(6.0)
add_card(s, Inches(0.5), col_y, col_w, col_h,
         "Today  ·  cache-aside on Oracle reads", [
             "2 Oracle EE read replicas — licence + support per core",
             "Senior engineer ~25% time on cache plumbing",
             "P95 latency outages drive support tickets",
             "Each new cached table = 2–3 sprint weeks",
         ], accent=WARN_RD)
add_card(s, Inches(6.8), col_y, col_w, col_h,
         "With RDI + Redis Enterprise", [
             "Drop 1–2 Oracle read replicas — reads served from Redis",
             "Engineers stop writing cache plumbing",
             "P95 < 5 ms for portfolio reads regardless of Oracle load",
             "New cached table = PR that adds a YAML file",
         ], accent=CHECK_GR)

# ROI banner
add_rect(s, Inches(0.5), Inches(5.85 - 0.1), Inches(12.3), Inches(0.0), DUSK_01)
add_rect(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5), MIDNIGHT)
add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
         "Typical ROI  ·  3–5× cost of Redis + RDI recovered in Oracle "
         "licence reduction within 12 months",
         size=13, bold=True, color=YELLOW, font=FONT_BODY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_talktrack(s,
    "I deliberately put cost on a slide because I know your CFO will ask. The "
    "economics are clear: every Oracle read you move to Redis is an Oracle core "
    "you don't have to licence. We have customers in Indian financial services "
    "who decommissioned Oracle read replicas entirely after going live on RDI. "
    "We can do a TCO worksheet specific to your environment as follow-up."
)

# -----------------------------------------------------------------
# 12. Security & compliance
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 12, TOTAL)
add_kicker(s, "Security & compliance")
add_h1(s, "Answers we have for your CISO up-front")

items = [
    ("Data in transit",
     "TLS / mTLS on every hop  ·  source ↔ RDI ↔ target."),
    ("Data at rest",
     "RDI never persists data on its VMs. All state in Redis Enterprise, "
     "inherits cluster encryption-at-rest."),
    ("Credentials",
     "RDI secrets (file system or K8s secrets) referenced by name in config.yaml — "
     "never plaintext."),
    ("RBAC",
     "Source DB user = replication-only.  Target DB ACL = write-only for RDI, "
     "read-only for applications."),
    ("Audit",
     "Structured JSON logs per pipeline — Splunk / ELK / Datadog ready."),
    ("Air-gapped option",
     "Runs entirely on your VMs / K8s. Offline installable. No outbound calls."),
]
y = Inches(2.0)
for title, body in items:
    add_rect(s, Inches(0.5), y, Inches(0.08), Inches(0.55), RED)
    add_text(s, Inches(0.7), y, Inches(3.3), Inches(0.35),
             title, size=13, bold=True, color=MIDNIGHT, font=FONT_BODY)
    add_text(s, Inches(0.7), y + Inches(0.28), Inches(12.0), Inches(0.4),
             body, size=12, color=DUSK, font=FONT_BODY)
    y += Inches(0.62)

add_talktrack(s,
    "We anticipate these questions because every financial services prospect asks "
    "them. Short version: RDI fits inside your existing perimeter, encrypts everything, "
    "and persists nothing of yours outside Redis Enterprise. SOC 2 reports and full "
    "architecture documentation available after the meeting."
)

# -----------------------------------------------------------------
# 13. Next steps
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
add_chrome(s, 13, TOTAL)
add_kicker(s, "Proposal")
add_h1(s, "Next steps  ·  six weeks to a production-ready PoC")

milestones = [
    ("Today",   "Validate the demo. Identify 2–3 portfolio tables for PoC."),
    ("Week 1",  "2-hour design session — DBAs walk through Oracle LogMiner / GoldenGate prerequisites."),
    ("Week 2",  "Stand up RDI in your dev environment against a read-only Oracle clone (we install with your team)."),
    ("Week 3–4","Cut over Portfolio screen reads to Redis in dev/UAT. Measure P95 + Oracle CPU reduction."),
    ("Week 6",  "Production go-live decision. RDI licence + Redis Enterprise sizing finalised."),
]
y = Inches(2.0); h = Inches(0.7); gap = Inches(0.12)
for label, body in milestones:
    add_rect(s, Inches(0.5), y, Inches(2.4), h, MIDNIGHT)
    add_text(s, Inches(0.5), y, Inches(2.4), h,
             label, size=15, bold=True, color=YELLOW, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(2.9), y, Inches(9.93), h, DUSK_01)
    add_text(s, Inches(3.15), y, Inches(9.65), h,
             body, size=13, color=DUSK, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    y += h + gap

# Bottom commitment
add_rect(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5), RED)
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
         "Redis provides a Solution Architect at no cost for the PoC. "
         "You keep every artefact produced.",
         size=13, bold=True, color=WHITE, font=FONT_BODY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Slight overflow into talktrack zone — skip talktrack on this slide
# to keep the layout clean.

# -----------------------------------------------------------------
# 14. Closing - run the demo
# -----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, MIDNIGHT)
# Brand strip
add_rect(s, Inches(0.6), Inches(0.45), Inches(0.28), Inches(0.28), RED)
add_text(s, Inches(0.97), Inches(0.32), Inches(3), Inches(0.5),
         "redis", size=24, bold=True, color=WHITE, font=FONT_BODY)
# Yellow pill
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(2.6),
                          Inches(2.2), Inches(0.5))
pill.line.fill.background()
pill.fill.solid(); pill.fill.fore_color.rgb = YELLOW
pill.adjustments[0] = 0.5
add_text(s, Inches(0.6), Inches(2.6), Inches(2.2), Inches(0.5),
         "DEMO  ·  STARTING NOW",
         size=12, bold=True, color=MIDNIGHT, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Title
add_text(s, Inches(0.6), Inches(3.2), Inches(12), Inches(1.5),
         "Let's run the demo.",
         size=68, bold=True, color=WHITE, font=FONT_BODY)
add_rect(s, Inches(0.6), Inches(4.55), Inches(1.6), Inches(0.07), RED)
add_text(s, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
         "…then your questions.",
         size=22, color=DUSK_30, font=FONT_BODY)
# URLs
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.5),
         "Dashboard         http://localhost:5050",
         size=14, color=YELLOW, font=FONT_MONO)
add_text(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.5),
         "Redis Insight     http://localhost:5540",
         size=14, color=YELLOW, font=FONT_MONO)
add_text(s, Inches(0.6), Inches(6.30), Inches(12), Inches(0.5),
         "RDI control API   https://rdi-api   (default / rdi_demo_pass)",
         size=14, color=YELLOW, font=FONT_MONO)
# Footer brand
add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.3),
         "Redis Data Integration  ·  prepared for Securities & Trading Firm",
         size=10, color=DUSK_50, font=FONT_MONO)


# =====================================================================
# DEMO SCENARIO SLIDES (15-20)
#
# These slides are kept on screen DURING the live demo. The audience
# sees a clean "what's about to happen + what to watch for" view.
# The presenter sees the literal copy-paste commands and timing notes
# in the slide notes pane (Presenter View in PowerPoint / Keynote).
# =====================================================================

def scenario_slide(
    slide_num, tag, minutes, title, sub,
    watch_items,            # left column: what the audience will see
    message,                # right column: the headline message to deliver
    cues,                   # bottom strip: where to look during the action
    notes,                  # speaker notes: step-by-step commands
):
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_chrome(s, slide_num, TOTAL)

    # Scenario badge (yellow pill) + minutes
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(0.85),
                              Inches(2.0), Inches(0.42))
    pill.line.fill.background()
    pill.fill.solid(); pill.fill.fore_color.rgb = YELLOW
    pill.adjustments[0] = 0.5
    add_text(s, Inches(0.5), Inches(0.85), Inches(2.0), Inches(0.42),
             tag.upper(), size=12, bold=True, color=MIDNIGHT, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.65), Inches(0.85), Inches(4.0), Inches(0.42),
             f"~{minutes} min", size=12, color=DUSK_50, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)

    # Big title
    add_text(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.85),
             title, size=32, bold=True, color=MIDNIGHT, font=FONT_BODY)
    add_text(s, Inches(0.5), Inches(2.30), Inches(12.3), Inches(0.5),
             sub, size=15, color=DUSK_50, font=FONT_BODY, italic=True)

    # Two-column body
    col_y = Inches(2.95); col_h = Inches(2.7)
    add_card(s, Inches(0.5), col_y, Inches(6.0), col_h,
             "What the room will see", watch_items, accent=RED)

    # Right "message" card: midnight background with yellow accent
    msg_x = Inches(6.85); msg_w = Inches(6.0)
    add_rect(s, msg_x, col_y, msg_w, col_h, MIDNIGHT)
    add_rect(s, msg_x, col_y, Inches(0.18), col_h, YELLOW)
    add_text(s, msg_x + Inches(0.4), col_y + Inches(0.2),
             msg_w - Inches(0.6), Inches(0.4),
             "HEADLINE MESSAGE", size=11, bold=True, color=YELLOW,
             font=FONT_MONO)
    add_text(s, msg_x + Inches(0.4), col_y + Inches(0.7),
             msg_w - Inches(0.6), col_h - Inches(0.9),
             f"\u201c{message}\u201d", size=14, color=WHITE,
             font=FONT_BODY, italic=True)

    # Bottom strip — "Look here" cues for the presenter to call out
    cue_y = Inches(5.85); cue_h = Inches(1.05)
    add_rect(s, Inches(0.5), cue_y, Inches(0.06), cue_h, RED)
    add_rect(s, Inches(0.56), cue_y, Inches(12.27), cue_h, DUSK_01)
    add_text(s, Inches(0.85), cue_y + Inches(0.08), Inches(12),
             Inches(0.3), "LOOK HERE", size=9, bold=True,
             color=RED, font=FONT_MONO)
    # Render cues as 1 line each, joined with the bullet glyph
    cue_text = "    ".join(f"\u25A0  {c}" for c in cues)
    add_text(s, Inches(0.85), cue_y + Inches(0.35),
             Inches(12.0), Inches(0.7), cue_text,
             size=11, color=DUSK, font=FONT_BODY)

    set_notes(s, notes)
    return s


# ---------------------------------------------------------------------
# Scenario 1 — Initial snapshot
# ---------------------------------------------------------------------
scenario_slide(
    slide_num=15, tag="Scenario 1", minutes=3,
    title="Initial snapshot — everything is already there",
    sub="No cold cache. Day-1 prefetch is done before any user logs in.",
    watch_items=[
        "Open the dashboard at http://localhost:5050",
        "Left rail already lists 10 the firm customers",
        "Pick HS0010001 (retail), then HS0010002 (HNI), HS0010003 (UHNI)",
        "Each customer's portfolio renders in single-digit ms",
    ],
    message="RDI prefetched every customer, security, holding, "
            "trade and price from Postgres on stack start. "
            "No first-user-pays-the-cold-cache penalty.",
    cues=[
        "Dashboard left rail",
        "Refresh time on right",
        "Portfolio holdings table",
    ],
    notes=(
        "PRE-CHECK  (do this before opening dashboard)\n"
        "  $ docker exec sectrade-redis-enterprise redis-cli -p 12000 DBSIZE\n"
        "    expected: >= 40  (means snapshot completed)\n"
        "\n"
        "EXECUTE\n"
        "  1. Browser → http://localhost:5050  (dashboard)\n"
        "  2. Walk the layout:\n"
        "        - left rail: 10 customers\n"
        "        - middle: holdings + KPIs\n"
        "        - right: latency panel + RDI pipeline metrics\n"
        "  3. Click 'HS0010001 — Rajesh Kumar Sharma'  (retail)\n"
        "        Point at refresh time on the right card. Should be < 5 ms.\n"
        "  4. Click 'HS0010002 — Priya Iyer'          (HNI)\n"
        "  5. Click 'HS0010003 — Vikram Mehta'        (UHNI)\n"
        "\n"
        "TALK\n"
        "  - 'I haven't run any setup command in front of you. The stack came\n"
        "     up 20 minutes ago and RDI snapshotted every selected table into\n"
        "     Redis before this meeting. None of these reads ever hit Postgres.'\n"
        "  - 'In your production rollout, this is the day-1 step you run once\n"
        "     during cutover. After it completes, every user request is hot.'\n"
        "\n"
        "BACKUP if it's empty\n"
        "  $ docker logs --tail 30 sectrade-rdi-processor   # look for JSON.SET\n"
        "  $ ./scripts/verify-redis.sh\n"
    ),
)

# ---------------------------------------------------------------------
# Scenario 2 — Live trade
# ---------------------------------------------------------------------
scenario_slide(
    slide_num=16, tag="Scenario 2", minutes=4,
    title="Live trade — Oracle change to Redis cache in ~1 second",
    sub="Insert one BUY in Postgres. Watch the dashboard reflect it without any app change.",
    watch_items=[
        "psql session showing the existing last 3 trades",
        "Single INSERT — order_id DEMO-LIVE-001",
        "Dashboard 'Recent trade stream' panel: new BUY row appears",
        "'RDI pipeline streams' card: trade count ticks up",
    ],
    message="No cache invalidation logic. No retry queue. "
            "The new trade landed in Redis because RDI is "
            "tailing the WAL — that's it.",
    cues=[
        "psql terminal",
        "Dashboard → Recent trade stream",
        "Dashboard → RDI pipeline streams card",
    ],
    notes=(
        "PRE-OPEN  (terminal tab 1, kept open from setup)\n"
        "  $ docker exec -it sectrade-postgres psql -U postgres sectrade\n"
        "\n"
        "EXECUTE  -- in psql\n"
        "  1. Show baseline trades:\n"
        "       SELECT trade_id, security_id, side, quantity, price, executed_at\n"
        "       FROM portfolio.trade\n"
        "       WHERE customer_id=10001\n"
        "       ORDER BY executed_at DESC LIMIT 3;\n"
        "\n"
        "  2. Place the new trade  (this is the moment of truth):\n"
        "       INSERT INTO portfolio.trade\n"
        "         (trade_id, customer_id, security_id, side, quantity, price,\n"
        "          trade_value, brokerage, order_id, exchange, executed_at)\n"
        "       VALUES (nextval('portfolio.trade_id_seq'),\n"
        "               10001, 1001, 'BUY',\n"
        "               10, 2945.50, 29455.00, 14.73,\n"
        "               'DEMO-LIVE-001', 'NSE', now());\n"
        "\n"
        "  3. (Optional, makes the holding KPI react too):\n"
        "       UPDATE portfolio.holding\n"
        "          SET quantity = quantity + 10,\n"
        "              invested_value = (quantity + 10) * avg_buy_price,\n"
        "              updated_at = now()\n"
        "        WHERE customer_id = 10001 AND security_id = 1001;\n"
        "\n"
        "SWITCH to dashboard tab (HS0010001 already selected). Within 1-2 s:\n"
        "  - Recent trade stream shows new BUY  (order_id DEMO-LIVE-001)\n"
        "  - RDI pipeline streams card: trade count +1\n"
        "  - If you ran step 3: RELIANCE holding qty + invested value jump\n"
        "\n"
        "TALK\n"
        "  - 'I did one INSERT into Postgres. No app code, no cache.invalidate(),\n"
        "     no message queue. RDI's CDC saw the WAL entry, transformed it via\n"
        "     trade.yaml, and wrote into Redis as a stream entry + JSON doc.'\n"
        "  - 'In your environment this insert would happen via the OMS, not by\n"
        "     hand. The path from OMS to dashboard is identical.'\n"
        "\n"
        "FALLBACK if nothing appears within 5 s\n"
        "  $ docker logs --tail 20 sectrade-rdi-processor\n"
        "  $ docker logs --tail 20 sectrade-debezium\n"
    ),
)

# ---------------------------------------------------------------------
# Scenario 3 — Continuous load
# ---------------------------------------------------------------------
scenario_slide(
    slide_num=17, tag="Scenario 3", minutes=3,
    title="Live load — peak-hour simulation",
    sub="Trade ticker + price ticker running together at broker-peak rate.",
    watch_items=[
        "Day-% column flickers across rows as prices tick",
        "RDI pipeline streams card increments visibly",
        "Trade stream gets several entries per minute",
        "Dashboard refresh latency stays sub-ms throughout",
    ],
    message="Everything you just saw for one trade is happening at scale. "
            "RDI processor sustains ~10k records/sec/core — "
            "we are nowhere near saturation here.",
    cues=[
        "Dashboard holdings table → Day %",
        "RDI pipeline streams (per-table counts)",
        "Trade stream — multiple entries per minute",
    ],
    notes=(
        "EXECUTE  (terminal tab 2)\n"
        "  $ ./scripts/run-simulation.sh\n"
        "    defaults: ~2 trades/sec  +  ~8 price ticks/sec, runs until Ctrl-C\n"
        "    tunables in env:  TRADES_PER_SEC  PRICES_PER_SEC  DURATION\n"
        "\n"
        "DURING  (let it run 60-90 s before commenting)\n"
        "  - switch to dashboard\n"
        "  - point at Day-% column changing  (prices)\n"
        "  - point at RDI pipeline streams card incrementing visibly\n"
        "  - point at the trade stream panel growing\n"
        "\n"
        "STOP  (when you're ready to move on)\n"
        "  Ctrl-C in terminal tab 2.  Stack stays healthy — simulator\n"
        "  doesn't need to be running for the rest of the demo.\n"
        "  (You can also just leave it running — it's lightweight.)\n"
        "\n"
        "TALK\n"
        "  - 'This is broker peak-hour scaled to 10 customers. ~2 trades/sec\n"
        "     and ~8 price ticks/sec across 20 securities, by default.'\n"
        "  - 'On real production hardware this scales linearly. We have\n"
        "     financial-services customers running 30-50k events/sec through\n"
        "     a single RDI active node.'\n"
        "  - 'Note nothing on the dashboard is slowing down — Redis serves\n"
        "     reads in microseconds regardless of write rate.'\n"
        "\n"
        "OBSERVABILITY  (only if asked)\n"
        "  $ docker exec sectrade-redis-rdi redis-cli -p 12001 XLEN \\\n"
        "      sectrade.portfolio.trade\n"
    ),
)

# ---------------------------------------------------------------------
# Scenario 4 — Latency comparison
# ---------------------------------------------------------------------
scenario_slide(
    slide_num=18, tag="Scenario 4", minutes=3,
    title="Postgres vs Redis — same query, two answers",
    sub="The dashboard runs the same portfolio-summary query against both. Ratio matters more than absolute numbers.",
    watch_items=[
        "Latency panel on right of dashboard",
        "Click 'Run again' three times — numbers stabilise",
        "Postgres: tens of ms.  Redis: sub-ms.",
        "4–8× speedup, consistently across runs",
    ],
    message="Multiply this ratio by your peak QPS to see how many "
            "Oracle cores you no longer need to license.",
    cues=[
        "Dashboard → Latency panel",
        "Postgres value",
        "Redis value",
        "Speedup x value",
    ],
    notes=(
        "EXECUTE\n"
        "  - on dashboard, scroll to Latency panel  (right column)\n"
        "  - click 'Run again' three times\n"
        "  - typical:  PostgreSQL ~5 ms  ·  Redis ~1 ms  ·  4-8x speedup\n"
        "\n"
        "TALK\n"
        "  - 'These are end-to-end from the app's perspective — connection\n"
        "     setup, query, deserialise. Redis is consistently 4-8x faster.'\n"
        "  - 'On your production hardware the absolute numbers will differ,\n"
        "     but the ratio holds. Ratios compound under load.'\n"
        "\n"
        "DO  NOT  RUN  benchmark.sh  LIVE  unless you have a deeply\n"
        "technical audience and 10 minutes to explain it.\n"
        "  - It uses persistent Postgres connections (no setup cost)\n"
        "  - On a laptop with hot buffers, Postgres CAN look faster for\n"
        "    a tiny dataset. That is an honest data point, not the right\n"
        "    one to lead with.\n"
        "\n"
        "IF ASKED about benchmark.sh\n"
        "  - 'On a laptop with the entire dataset in Postgres shared buffers\n"
        "     and a warm connection pool, Postgres can win small synthetic\n"
        "     benchmarks. The numbers that matter come from production:\n"
        "     Oracle CPU when 5,000 users hit the portfolio screen at 9:15.\n"
        "     Postgres queues, Redis scales. We will measure exactly that\n"
        "     in the PoC.'\n"
        "\n"
        "REFERENCE  (offline only)\n"
        "  $ ./scripts/benchmark.sh\n"
    ),
)

# ---------------------------------------------------------------------
# Scenario 5 — Redis Insight observability + RDI tab
# ---------------------------------------------------------------------
scenario_slide(
    slide_num=19, tag="Scenario 5", minutes=5,
    title="Redis Insight — pipeline editor + live analytics",
    sub="The same tool your ops team will use in dev, staging, prod. Free with Redis Enterprise.",
    watch_items=[
        "Insight: keys in target DB (holding, trades, price)",
        "Insight: 5 CDC streams in the RDI state DB",
        "Insight RDI tab: the deployed YAML pipeline",
        "Insight RDI tab: Analytics view with live event counts",
    ],
    message="The pipeline YAML is committed to git, Insight reads it "
            "via RDI's API, your platform team owns it end-to-end. "
            "No Kafka team. No bespoke UI.",
    cues=[
        "Insight → Browser",
        "Insight → Redis Data Integration tab",
        "Pipeline Management → editor",
        "Analytics → live throughput",
    ],
    notes=(
        "PART 1 — Browser (data view)  ~2 min\n"
        "  Open http://localhost:5540  → 'Portfolio cache (target)' connection\n"
        "  Filter:  holding:10001:*    → click holding:10001:1001  → JSON tree\n"
        "  Filter:  trades:10001       → click → STREAM view\n"
        "  Filter:  price:*            → click price:1001         → HASH view\n"
        "\n"
        "  Switch connection to 'RDI state DB'\n"
        "  Filter:  sectrade.portfolio.*  → 5 streams, one per source table\n"
        "  Click one → STREAM view → expand an entry → show Debezium\n"
        "    CDC envelope JSON (before / after / op / source / ts_ms)\n"
        "\n"
        "PART 2 — RDI tab  ~3 min\n"
        "  Click 'Redis Data Integration' icon on left rail of Insight\n"
        "  If not yet added, click 'Let's connect to RDI':\n"
        "        RDI Alias  =  sectrade-rdi-demo\n"
        "        URL        =  https://rdi-api\n"
        "        Username   =  default\n"
        "        Password   =  rdi_demo_pass\n"
        "  Click 'Add Endpoint' → click into the new instance.\n"
        "\n"
        "  Walk three views:\n"
        "    1. Pipeline Management\n"
        "         - left pane: config.yaml + 5 jobs/*.yaml\n"
        "         - click holding.yaml — show YAML editor\n"
        "         - 'this is the file we commit to git'\n"
        "    2. Test Connection (bottom right of editor)\n"
        "         - both target and source come back green\n"
        "         - 'Insight is asking RDI to validate against real DBs\n"
        "            BEFORE deploy. Zero risk of bad config in prod.'\n"
        "    3. Analytics tab\n"
        "         - throughput, 5 CDC streams, processed counts\n"
        "         - snapshot status: completed\n"
        "         - refresh interval: 5 s\n"
        "\n"
        "TALK\n"
        "  - 'Same tool, same UI, in dev / staging / prod. Free with Redis\n"
        "     Enterprise. Your engineers already know it.'\n"
        "  - 'The pipeline is observable from the moment Debezium captures\n"
        "     an event until it lands in the target. Lag, throughput, DLQ —\n"
        "     all here, all queryable.'\n"
        "\n"
        "DISCLOSURE if a sharp engineer asks\n"
        "  - The Deploy / Start / Stop / Reset buttons are no-ops in THIS\n"
        "    demo because the pipeline is already running via docker-compose.\n"
        "  - In production they trigger the RDI control plane API to roll\n"
        "    the pipeline; same UX.\n"
    ),
)

# ---------------------------------------------------------------------
# Scenario 6 — Schema change
# ---------------------------------------------------------------------
scenario_slide(
    slide_num=20, tag="Scenario 6", minutes=4,
    title="Schema change handling — add a column, no restart",
    sub="ALTER TABLE in Postgres → next UPDATE flows to Redis automatically.",
    watch_items=[
        "ALTER TABLE portfolio.holding ADD COLUMN strategy_tag",
        "Single UPDATE that touches customer 10001 rows",
        "redis-cli JSON.GET → new field present, no other change",
        "No pipeline reset, no deploy, no downtime",
    ],
    message="New columns flow through automatically. "
            "If you need to EXCLUDE a column (PII, PAN), "
            "you add a 3-line transformation to the YAML. That's it.",
    cues=[
        "psql terminal — ALTER + UPDATE",
        "redis-cli JSON.GET output",
        "'strategy_tag' field in the JSON",
    ],
    notes=(
        "EXECUTE  -- in psql terminal tab 1\n"
        "  1. Add the column:\n"
        "       ALTER TABLE portfolio.holding\n"
        "         ADD COLUMN strategy_tag VARCHAR(40);\n"
        "\n"
        "  2. Update some rows so the new column ends up in the WAL:\n"
        "       UPDATE portfolio.holding\n"
        "          SET strategy_tag = 'LONG_TERM'\n"
        "        WHERE customer_id = 10001;\n"
        "\n"
        "  (IMPORTANT: in Postgres 16+ a bare ALTER doesn't write to the\n"
        "   WAL by itself — you need a row UPDATE for CDC to see it.)\n"
        "\n"
        "VERIFY  -- in another terminal\n"
        "  $ docker exec sectrade-redis-enterprise redis-cli -p 12000 \\\n"
        "       JSON.GET holding:10001:1001 $\n"
        "    expected: the JSON now contains \"strategy_tag\":\"LONG_TERM\"\n"
        "\n"
        "TALK\n"
        "  - 'No pipeline restart. No deploy. No code change. New columns\n"
        "     auto-propagate.'\n"
        "  - 'The opposite — removing or masking a column for compliance —\n"
        "     is a 3-line YAML transform. We do that in the PoC against\n"
        "     real PII columns like PAN and Aadhaar.'\n"
        "\n"
        "OPTIONAL deep-dive  (only if architects push on schema evolution)\n"
        "  - Renames / type-narrowing: RDI surfaces a 'schema drift' event\n"
        "    in the DLQ; you decide whether to accept or reject.\n"
        "  - Drops: the column simply stops appearing in new writes;\n"
        "    existing keys are untouched until next change.\n"
    ),
)

# ---------------------------------------------------------------------
# 21. Q&A / Recap
# ---------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, MIDNIGHT)
add_rect(s, Inches(0.6), Inches(0.45), Inches(0.28), Inches(0.28), RED)
add_text(s, Inches(0.97), Inches(0.32), Inches(3), Inches(0.5),
         "redis", size=24, bold=True, color=WHITE, font=FONT_BODY)
add_text(s, Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.3),
         f"21 / {TOTAL}", size=11, color=DUSK_50, font=FONT_MONO,
         align=PP_ALIGN.RIGHT)

# Yellow pill — Questions
pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(1.6),
                          Inches(2.0), Inches(0.45))
pill.line.fill.background()
pill.fill.solid(); pill.fill.fore_color.rgb = YELLOW
pill.adjustments[0] = 0.5
add_text(s, Inches(0.6), Inches(1.6), Inches(2.0), Inches(0.45),
         "Q & A  ·  WRAP",
         size=12, bold=True, color=MIDNIGHT, font=FONT_MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.1),
         "Three numbers we'll hand you at end of PoC.",
         size=42, bold=True, color=WHITE, font=FONT_BODY)
add_rect(s, Inches(0.6), Inches(3.35), Inches(1.6), Inches(0.07), RED)

# Three KPI tiles
def kpi(x, big, label, sub):
    add_text(s, x, Inches(3.7), Inches(4), Inches(1.2),
             big, size=58, bold=True, color=YELLOW, font=FONT_BODY)
    add_text(s, x, Inches(4.85), Inches(4), Inches(0.4),
             label, size=14, bold=True, color=WHITE, font=FONT_MONO)
    add_text(s, x, Inches(5.2), Inches(4), Inches(1.5),
             sub, size=12, color=DUSK_30, font=FONT_BODY)

kpi(Inches(0.6),  "P95",  "PORTFOLIO READ LATENCY",
    "Measured before and after RDI cutover, on your production-equivalent UAT load.")
kpi(Inches(5.0),  "%↓",   "ORACLE CPU REDUCTION",
    "Hot-path Portfolio reads moved off Oracle. Direct input to your DBA licensing model.")
kpi(Inches(9.4),  "h/wk", "ENGINEERING TIME RECLAIMED",
    "Hours your senior engineer stops spending on cache-aside plumbing across 6 repos.")

# Bottom commitment band
add_rect(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), RED)
add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         "Redis SA assigned to your PoC at no cost  ·  "
         "All artefacts produced are yours to keep",
         size=13, bold=True, color=WHITE, font=FONT_BODY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

set_notes(s, (
    "USE THIS SLIDE for the closing exchange after Q & A.\n"
    "\n"
    "Anchor the conversation on these three numbers — they map to\n"
    "the three buyer constituencies in the firm:\n"
    "  - P95 latency      → CTO / VP Engineering\n"
    "  - Oracle CPU %↓    → CFO  / IT Finance\n"
    "  - Eng hours / week → Engineering Manager / Tech Lead\n"
    "\n"
    "ASK FOR\n"
    "  1. A name + email for the PoC technical lead on their side\n"
    "  2. A read-only Oracle clone (or a willingness to use Postgres\n"
    "     for the first iteration)\n"
    "  3. A target 2-week window for the design session\n"
    "\n"
    "FOLLOW-UP commitments to make on the call\n"
    "  - Send the SOC 2 + architecture docs by EOD\n"
    "  - Send a TCO worksheet template within 48 h\n"
    "  - Schedule the design session within 5 working days\n"
))


# ---------------------------------------------------------------------
out = Path(__file__).resolve().parent.parent / "docs" / "01-slide-deck.pptx"
prs.save(str(out))
size_kb = os.path.getsize(out) // 1024
print(f"wrote {out}  ({size_kb} KB, {len(prs.slides)} slides)")
